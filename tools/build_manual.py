"""Assemble docs/screenshots/*.png into a PowerPoint manual.

Output:
    docs/AIVEX_Processing_Tool_v2.0_Manual.pptx

Runs against whatever PNGs are sitting in docs/screenshots/ — re-run
capture_screens.py first if the UI changed.
"""

from __future__ import annotations

import sys
from dataclasses import dataclass
from pathlib import Path

# Make the project root importable.
_PROJECT_ROOT = Path(__file__).resolve().parent.parent
if str(_PROJECT_ROOT) not in sys.path:
    sys.path.insert(0, str(_PROJECT_ROOT))

if hasattr(sys.stdout, "reconfigure"):
    try:
        sys.stdout.reconfigure(encoding="utf-8", errors="replace")
    except Exception:
        pass

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Emu, Inches, Pt

from apt.brand import APP_NAME, APP_VERSION


# AIVEX brand palette
ORANGE       = RGBColor(0xFF, 0x70, 0x29)
ORANGE_DARK  = RGBColor(0xC8, 0x50, 0x1A)
BLACK        = RGBColor(0x0B, 0x0B, 0x0E)
PANEL        = RGBColor(0x15, 0x16, 0x1B)
PANEL_2      = RGBColor(0x1D, 0x1F, 0x26)
WHITE        = RGBColor(0xED, 0xED, 0xEF)
DIM          = RGBColor(0x9A, 0x9C, 0xA3)
GREEN        = RGBColor(0x33, 0xB6, 0x6B)
BLUE         = RGBColor(0x5B, 0xA9, 0xF5)

SCREENSHOT_DIR = _PROJECT_ROOT / "docs" / "screenshots"
OUTPUT_PPTX    = _PROJECT_ROOT / "docs" / f"AIVEX_Processing_Tool_v{APP_VERSION}_Manual.pptx"


@dataclass(frozen=True)
class Slide:
    title: str
    body: str             # markdown-ish bullet text (use \n for newlines, • for bullets)
    image: str | None     # filename in docs/screenshots/, or None for text-only
    caption: str = ""     # small grey line under the image


SLIDES: list[Slide] = [
    # ---- Cover ----------------------------------------------------
    Slide(
        title="__COVER__",
        body=f"AIVEX Processing Tool\nv{APP_VERSION} 매뉴얼\n\n"
             f"v1 monolith → v2.0 패키지화 + 노드 그래프 기반 Preprocessing 패널",
        image=None,
    ),

    # ---- Overview ------------------------------------------------
    Slide(
        title="한눈에 보기",
        body=(
            "좌측 sidebar에서 작업 선택 → 우측에서 처리\n"
            "총 11개 패널 (v1 10개 + Preprocessing 신규) · 4개 카테고리\n\n"
            "• Sorting   — Basic Sorting · NG Folder Sorting · NG Count\n"
            "• Copy      — Date-Based Copy · Image Format Copy · Simulation Foldering\n"
            "• Image Ops — Crop · Attach FOV · BMP to JPG · Preprocessing ⭐\n"
            "• Conversion — MIM to BMP"
        ),
        image="01_main_basic_sorting.png",
        caption="Basic Sorting 패널 — 모든 v1 패널은 같은 뼈대(Configuration form + Logs + progress + Start/Stop)를 따릅니다.",
    ),

    # ---- Preprocessing panel intro -------------------------------
    Slide(
        title="🎨 Preprocessing 패널 — v2.0의 메인 신규 기능",
        body=(
            "이미지 전처리를 시각적 노드 그래프로 구성하고, 여러 이미지에 한 번에 적용하고,\n"
            ".apt.json으로 저장해 다음에 그대로 재사용.\n\n"
            "핵심 컨셉:\n"
            "• Origin — 로드한 이미지가 들어오는 시작 노드 (1개, 고정)\n"
            "• Op 노드 — 31개 전처리 op 중 선택해서 그래프에 배치\n"
            "• 연결 — output port (오른쪽 초록) → input port (왼쪽 파랑) 드래그\n"
            "• Leaf — 다른 노드의 input이 아닌 끝단. Save Outputs 시 leaf만 저장\n"
            "• Active 이미지 — 인스펙터 Active 탭이 보여주는 기준 이미지"
        ),
        image="02_preprocessing_empty.png",
        caption="Origin만 있는 빈 캔버스. 좌측 Operations · 가운데 Pipeline Graph · 우측 Properties/Parameters/Preview 3-pane 레이아웃.",
    ),

    # ---- Workflow step 1 -----------------------------------------
    Slide(
        title="워크플로 1단계 — 이미지 로드",
        body=(
            "Load Images... / Add Images... 또는 Load Samples 버튼\n"
            "JPG / PNG / BMP 다중 선택 (MIM은 MIM to BMP로 먼저 변환)\n\n"
            "로드되면:\n"
            "• 캔버스 아래 Image Strip에 thumbnail 표시\n"
            "• 첫 이미지가 자동 active 지정\n"
            "• Strip 클릭으로 active 변경 · A / D 키로 빠른 전환\n"
            "• 이미지 2장 이상이면 기본 탭이 All Images로 자동 전환"
        ),
        image="03_preprocessing_samples_loaded.png",
        caption="번들된 sample 이미지 2장이 Image Strip에 추가된 상태.",
    ),

    # ---- Workflow step 2 -----------------------------------------
    Slide(
        title="워크플로 2단계 — 노드 추가 & 연결",
        body=(
            "좌측 Operations 검색창에 op 이름 입력 (blur, edge, threshold…)\n"
            "카드를 더블클릭하면 캔버스에 노드 추가\n\n"
            "연결: 노드의 오른쪽 초록 port → 다른 노드의 왼쪽 파랑 port로 드래그\n"
            "• Origin은 여러 갈래로 fan-out 가능\n"
            "• Combine 노드 (Blend/Add/…)는 input port 2개\n"
            "• edge는 source 노드의 카테고리 컬러를 따라 그려짐\n\n"
            "Auto-fork: 단일-input op의 이미 차있는 port에 두 번째 edge를 떨어뜨리면\n"
            "destination 노드가 자동 복제됩니다 (같은 params, 이후 독립 편집)"
        ),
        image="04_preprocessing_one_node.png",
        caption="Origin → Grayscale 첫 연결. 노드 카드 우상단의 status pip와 실행 시간 라벨이 보입니다.",
    ),

    # ---- Workflow step 3 -----------------------------------------
    Slide(
        title="워크플로 3단계 — 파라미터 + 실시간 미리보기",
        body=(
            "노드 클릭 → 우측 Parameters에 자동 폼 생성 (spinbox / checkbox / combo)\n"
            "값을 바꾸면 50ms debounce로 Active 탭이 재계산\n\n"
            "노드 상태 (Cognex 스타일):\n"
            "• 🟢 success — 방금 실행 성공\n"
            "• 🔵 cached — 캐시 재사용\n"
            "• 🔴 error — 실패 (호버 툴팁에 메시지)\n"
            "• ⚫ idle — 아직 실행 안 함\n\n"
            "Properties 패널: Name / Type / Status / Time / I-O / Output Shape 한눈에"
        ),
        image="05_preprocessing_full_pipeline_active.png",
        caption="균열 검출 레시피: Origin → Grayscale → Crop → Window Stretch → Resize Smooth. 각 노드의 실행 시간이 표시됩니다.",
    ),

    # ---- All images grid ----------------------------------------
    Slide(
        title="다중 이미지 — All Images 탭",
        body=(
            "선택한 노드를 모든 이미지에 적용한 결과를 3열 grid로 표시.\n"
            "한 번 파이프라인 짜면 N장에 동일하게 적용해 일관성 확인.\n\n"
            "성능: Grid는 120ms debounce로 비동기 계산.\n"
            "이미지 10장 넘으면 0.5초 정도 갱신 delay 가능.\n\n"
            "Active 탭과 자동 전환:\n"
            "• 이미지 1장 — 기본 탭은 Active\n"
            "• 이미지 2장+ — 기본 탭은 All Images"
        ),
        image="06_preprocessing_all_images_grid.png",
        caption="이미지 2장에 동일 파이프라인을 적용한 결과를 grid로 비교.",
    ),

    # ---- Fan-out ------------------------------------------------
    Slide(
        title="Fan-out & Auto-fork",
        body=(
            "한 source가 여러 destination에 연결 가능 (fan-out).\n"
            "균열 검출에서 윈도우를 슬라이딩할 때 자주 씀:\n"
            "Crop output → Window Stretch(80,180) / (120,200) / (160,240)…\n\n"
            "Auto-fork:\n"
            "단일-input op의 이미 차있는 port에 두 번째 edge를 끌면\n"
            "destination 노드가 자동 복제됩니다 (같은 params, 이후 독립 편집).\n"
            "Ctrl+D로 명시적 복제와는 다르게, edge를 끄는 자연스러운 동작 중에 발동."
        ),
        image="07_preprocessing_fanout.png",
        caption="Crop 결과를 3개의 Window Stretch로 fan-out. 각 윈도우 범위가 다른 결과를 동시에 생성.",
    ),

    # ---- Job file + Export --------------------------------------
    Slide(
        title="Job 파일 + Save Outputs",
        body=(
            "Job 파일 (.apt.json) — 그래프 전체를 JSON으로 저장\n"
            "• Save Job... → nodes, params, connections, positions 저장\n"
            "• Load Job... → 다른 이미지 셋에 그대로 적용\n"
            "• 팀 표준 recipe 공유: crack_detection_v1.apt.json 식\n\n"
            "Save Outputs... — (이미지 × leaf 노드) 모든 조합을 풀해상도 PNG로 저장\n"
            "• 파일명: <원본 stem>__<node-id>.png\n"
            "• Preview는 720px downscale, Export는 원본 해상도\n"
            "• 일부 노드 실패해도 나머지 저장 계속 진행"
        ),
        image=None,
    ),

    # ---- Operations catalog -------------------------------------
    Slide(
        title="31개 Operations 카탈로그",
        body=(
            "Geometry — Resize · Rotate · Flip · Crop (XYWH)\n"
            "Color — Grayscale · Invert · Brightness/Contrast · Gamma · Window Stretch\n"
            "Filter — Gaussian / Median / Bilateral / Box Blur · Unsharp Mask · Resize Smooth\n"
            "Threshold — Binary · Otsu · Adaptive (gaussian/mean)\n"
            "Edge — Canny · Sobel · Laplacian\n"
            "Morphology — Erode · Dilate · Open · Close\n"
            "Histogram — Equalize · CLAHE\n"
            "Combine (2-input) — Blend · Add · Subtract · Max · Min\n\n"
            "새 op 추가: apt/preprocessing/operations.py에 Operation(...) 한 줄 +\n"
            "함수 하나면 sidebar / parameter form / tests에 자동 반영"
        ),
        image=None,
    ),

    # ---- Shortcuts ----------------------------------------------
    Slide(
        title="단축키 전체",
        body=(
            "캔버스 — 이동\n"
            "• 노드 드래그 — snap-to-grid (Shift = 자유)  · Space + 드래그 — pan\n"
            "• 가운데 버튼 드래그 — pan  · 휠 — zoom  · F — Fit View  · Ctrl+0 — zoom 100%\n\n"
            "캔버스 — 선택\n"
            "• 클릭 단일 / Ctrl+클릭 추가 / 드래그 박스 선택 / Ctrl+A 전체 / Esc 해제\n\n"
            "캔버스 — 편집\n"
            "• Del/Backspace — 삭제  · Ctrl+D — duplicate  · ←↑↓→ nudge (Shift ×5)\n"
            "• 우클릭 노드 — Duplicate / Disconnect inputs / Delete 메뉴\n"
            "• Auto-Layout 버튼 — Origin → leaf 깊이 기준 자동 정렬\n\n"
            "이미지 탐색\n"
            "• A 이전 / D 다음 (텍스트 입력 포커스 있을 땐 자동 비활성)"
        ),
        image=None,
    ),

    # ---- 그 외 패널 -----------------------------------------------
    Slide(
        title="그 외 패널 — NG Count",
        body=(
            "NG 폴더의 Cam_*/Defect/* 항목 수 집계.\n"
            "결과는 Cam × Defect 카운트 표 + 상단 요약.\n\n"
            "Copy 버튼으로 클립보드에 tab-separated 복사 → Excel 바로 붙여넣기."
        ),
        image="08_ng_count.png",
        caption="NG Count 패널.",
    ),
    Slide(
        title="그 외 패널 — MIM to BMP",
        body=(
            "INI 파일을 편집해서 mim2color.exe 실행.\n"
            "Preprocessing 패널이 MIM을 직접 못 읽으니, MIM 데이터는 먼저 이 패널로 BMP 변환 필수.\n\n"
            "워크플로:\n"
            "1) Select INI file → 자동으로 PATH 섹션 초기화\n"
            "2) INI editor에서 Source/Target path 직접 입력\n"
            "3) Save → Start → mim2color.exe가 새 console로 실행"
        ),
        image="09_mim_to_bmp.png",
        caption="MIM to BMP 패널 — INI editor.",
    ),
    Slide(
        title="그 외 패널 — Crop",
        body=(
            "이미지 트리에서 FOV(optional)에 해당하는 파일을 일괄 crop.\n"
            "ltrb / xywh 좌표 모드 선택.\n\n"
            "BMP + JSON pair 처리:\n"
            "• 동명의 .json (LabelMe 포맷)이 있으면 라벨 좌표도 같이 보정\n"
            "• _draw.bmp debug 이미지도 생성"
        ),
        image="10_crop_panel.png",
        caption="Crop 패널.",
    ),

    # ---- Installer / release ----------------------------------
    Slide(
        title="빌드 · 배포 · 릴리즈",
        body=(
            "단일 인스톨러 (APT_Setup_v<버전>.exe) 한 파일로 배포\n\n"
            "사용자 PC에 설치되는 4 셋트:\n"
            "• APT.exe — 메인 실행 파일\n"
            "• _internal\\ — Python 런타임 + Qt DLL + 번들 리소스 (sample 이미지 포함)\n"
            "• mim2color.exe — MIM → BMP 외부 변환기\n"
            "• mim_converter_config.ini — 기본 INI (사용자 편집 보존)\n\n"
            "빌드: installer\\build.ps1 한 번 → ~4분\n"
            "버전 관리: apt/brand.py의 APP_VERSION 한 줄이 single source of truth\n"
            "Inno Setup 6 필요 (jrsoftware.org/isdl.php)"
        ),
        image=None,
    ),

    # ---- 마무리 -----------------------------------------------
    Slide(
        title="요약",
        body=(
            "v2.0의 핵심 가치\n\n"
            "✅ 코드 모듈화 — 6,415줄 monolith → apt/ 패키지 (50+ 파일)\n"
            "✅ AIVEX 브랜드 UI — 블랙/오렌지 sidebar + stacked pages\n"
            "✅ Preprocessing 패널 — 노드 그래프 + 31 op + 다중 이미지 + Job 파일\n"
            "✅ 단위 테스트 127개 (v1은 0개)\n"
            "✅ unhandled exception → error.log 자동 기록\n"
            "✅ 단일 인스톨러 배포 + 자동 업그레이드\n\n"
            "Repo: github.com/CVKim/AiV_ProcessingTool\n"
            "공유 인스톨러: \\\\192.168.10.74\\02.프로그램\\10. APT\\3. installer\\"
        ),
        image=None,
    ),
]


# -------------------------------------------------------------------------
# Slide builders
# -------------------------------------------------------------------------

def _add_fill(shape, color):
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_text(slide, left, top, width, height, text, *,
              size=18, bold=False, color=WHITE, align=PP_ALIGN.LEFT,
              anchor=MSO_ANCHOR.TOP, line_spacing=1.25):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    lines = text.split("\n")
    for i, line in enumerate(lines):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.alignment = align
        para.line_spacing = line_spacing
        run = para.add_run()
        run.text = line
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
        run.font.name = "Malgun Gothic"
    return tb


def _build_cover_slide(prs, slide_def):
    layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(layout)
    # Black background
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    _add_fill(bg, BLACK)
    bg.line.fill.background()

    # Big orange brand
    _add_text(
        slide, Inches(0.7), Inches(2.2), Inches(11.9), Inches(1.2),
        "AIVEX", size=72, bold=True, color=ORANGE,
    )
    _add_text(
        slide, Inches(0.75), Inches(3.3), Inches(11.9), Inches(0.6),
        "PROCESSING TOOL", size=20, bold=True, color=DIM,
    )
    # Subtitle
    body_lines = slide_def.body.split("\n")
    _add_text(
        slide, Inches(0.75), Inches(4.4), Inches(11.9), Inches(2.5),
        "\n".join(body_lines[1:]),  # skip the first line (we already showed AIVEX)
        size=18, color=WHITE,
    )

    # Bottom accent bar
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, prs.slide_height - Inches(0.18),
        prs.slide_width, Inches(0.18),
    )
    _add_fill(bar, ORANGE)
    bar.line.fill.background()


def _build_content_slide(prs, slide_def):
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)

    # Black background
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    _add_fill(bg, BLACK)
    bg.line.fill.background()

    # Top orange stripe
    stripe = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.08),
    )
    _add_fill(stripe, ORANGE)
    stripe.line.fill.background()

    # Title
    _add_text(
        slide, Inches(0.55), Inches(0.3), Inches(12.3), Inches(0.7),
        slide_def.title, size=28, bold=True, color=WHITE,
    )

    if slide_def.image and (SCREENSHOT_DIR / slide_def.image).is_file():
        # Image on the right; body on the left
        # Layout: 12-col grid. Left col = body (5 cols), right = image (7 cols)
        # Slide is 13.333 x 7.5 in. Available area below title: y from 1.1 to 7.3.
        body_left   = Inches(0.55)
        body_top    = Inches(1.15)
        body_width  = Inches(5.2)
        body_height = Inches(5.9)
        _add_text(
            slide, body_left, body_top, body_width, body_height,
            slide_def.body, size=13, color=WHITE, line_spacing=1.4,
        )

        img_path = str(SCREENSHOT_DIR / slide_def.image)
        img_left   = Inches(6.0)
        img_top    = Inches(1.15)
        img_width  = Inches(7.05)
        # Add image first to get its natural size, then resize.
        pic = slide.shapes.add_picture(
            img_path, img_left, img_top, width=img_width,
        )
        # Move down slightly if necessary so caption fits below
        caption_top = Inches(1.15) + pic.height + Inches(0.08)
        if slide_def.caption:
            _add_text(
                slide, img_left, caption_top, img_width, Inches(0.4),
                slide_def.caption, size=9, color=DIM, line_spacing=1.2,
            )
    else:
        # Text-only slide: wide body
        _add_text(
            slide, Inches(0.55), Inches(1.15), Inches(12.3), Inches(5.9),
            slide_def.body, size=16, color=WHITE, line_spacing=1.6,
        )

    # Footer
    _add_text(
        slide, Inches(0.55), Inches(7.05), Inches(12.3), Inches(0.35),
        f"AIVEX Processing Tool · v{APP_VERSION}",
        size=9, color=DIM, align=PP_ALIGN.RIGHT,
    )


def main() -> int:
    prs = Presentation()
    prs.slide_width  = Inches(13.333)   # 16:9
    prs.slide_height = Inches(7.5)

    for slide_def in SLIDES:
        if slide_def.title == "__COVER__":
            _build_cover_slide(prs, slide_def)
        else:
            _build_content_slide(prs, slide_def)

    OUTPUT_PPTX.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUTPUT_PPTX))
    print(f"Saved: {OUTPUT_PPTX}")
    print(f"Size:  {OUTPUT_PPTX.stat().st_size / 1024:.1f} KB")
    print(f"Slides: {len(SLIDES)}")
    return 0


if __name__ == "__main__":
    sys.exit(main())
